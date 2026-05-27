"""
Selva Digital — Black Friday 2026
Presentacion de Estrategia de Campana
16 slides · Tipo: custom · Audiencia: interna + stakeholders

Run:  python build_bf_2026.py
Out:  campana-black-friday-2026-06.pptx
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

# ============================================================
# DESIGN SYSTEM — Selva Digital tokens
# ============================================================
ACCENT       = RGBColor(0x2B, 0xB6, 0x73)   # #2BB673
BG_DARK      = RGBColor(0x0A, 0x0B, 0x0D)   # #0A0B0D
BG_LIGHT     = RGBColor(0xFA, 0xFA, 0xFA)   # #FAFAFA
SURFACE_1    = RGBColor(0x12, 0x13, 0x16)   # #121316
SURFACE_2    = RGBColor(0x1A, 0x1C, 0x20)   # #1A1C20
TEXT_PRIMARY = RGBColor(0x1C, 0x1C, 0x1F)   # #1C1C1F
TEXT_ON_DARK = RGBColor(0xFA, 0xFA, 0xFA)
TEXT_SOFT    = RGBColor(0x6B, 0x6B, 0x6E)
TEXT_DIM     = RGBColor(0x8E, 0x8E, 0x92)
CTA_TEXT     = RGBColor(0x06, 0x14, 0x0C)
ACCENT_TINT  = RGBColor(0xF0, 0xFB, 0xF6)

FONT_H = "Geist"
FONT_B = "Inter"
FONT_M = "JetBrains Mono"

SW   = Inches(13.333)
SH   = Inches(7.5)
MH   = Inches(0.833)
MV   = Inches(0.556)
SAFW = SW - 2 * MH
SAFH = SH - 2 * MV

S4  = Emu(45720)
S8  = Emu(91440)
S12 = Emu(137160)
S16 = Emu(182880)
S24 = Emu(273600)
S32 = Emu(365760)
S48 = Emu(548640)
S64 = Emu(731520)

# ============================================================
# HELPERS
# ============================================================

def new_prs():
    p = Presentation()
    p.slide_width  = SW
    p.slide_height = SH
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def tx(slide, text, x, y, w, h,
       fnt=FONT_B, sz=Pt(14), bold=False, italic=False,
       color=TEXT_PRIMARY, align=PP_ALIGN.LEFT):
    b  = slide.shapes.add_textbox(x, y, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text            = text
    r.font.name       = fnt
    r.font.size       = sz
    r.font.bold       = bold
    r.font.italic     = italic
    r.font.color.rgb  = color
    return b

def bar(slide, x, y, w, h=S4, color=ACCENT):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def rct(slide, x, y, w, h, fill, lc=None, lw=Emu(12700)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if lc:
        s.line.color.rgb = lc
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def pill(slide, text, x, y):
    """Terminal chip: dark bg, green border, mono font."""
    w = Emu(int(len(text) * 85000 + 300000))
    h = Emu(290000)
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = BG_DARK
    s.line.color.rgb = ACCENT
    s.line.width = Emu(12700)
    tf = s.text_frame
    tf.margin_left   = S8
    tf.margin_right  = S8
    tf.margin_top    = S4
    tf.margin_bottom = S4
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text            = text
    r.font.name       = FONT_M
    r.font.size       = Pt(10)
    r.font.color.rgb  = ACCENT
    return s

def divider_slide(prs, number, title_line1, title_line2=None):
    """Reusable section divider: dark surface, left accent bar, big heading."""
    s = blank(prs)
    bg(s, SURFACE_1)
    bar(s, Inches(0), SH * 0.18, S8, SH * 0.64)
    tx(s, number,
       MH + S24, SH * 0.24, Inches(2), Emu(380000),
       fnt=FONT_M, sz=Pt(11), color=ACCENT)
    title = title_line1 + ("\n" + title_line2 if title_line2 else "")
    tx(s, title,
       MH + S24, SH * 0.32, SAFW, Inches(2.5),
       fnt=FONT_H, sz=Pt(52), bold=True, color=TEXT_ON_DARK)
    return s

# ============================================================
# SLIDE BUILDERS
# ============================================================

def s01_cover(prs):
    s = blank(prs)
    bg(s, BG_DARK)

    # Accent vertical bar — right edge
    bar(s, SW - Inches(0.9), MV, S4, SH - 2 * MV)

    # Kicker chip
    pill(s, "ESTRATEGIA DE CAMPANA  2026", MH, MV)

    # Logo
    tx(s, "Selva Digital",
       MH, MV + Emu(420000), Inches(3.2), Emu(420000),
       fnt=FONT_H, sz=Pt(14), bold=True, color=ACCENT)

    # Main headline — 3 lines for rhythm
    tx(s, "Black Friday 2026:",
       MH, Inches(1.9), SAFW * 0.75, Inches(0.95),
       fnt=FONT_H, sz=Pt(50), bold=True, color=TEXT_ON_DARK)

    tx(s, "La semana en que las\nPyMEs compran.",
       MH, Inches(2.8), SAFW * 0.75, Inches(1.9),
       fnt=FONT_H, sz=Pt(50), bold=True, color=ACCENT)

    # Subtitle
    tx(s, "Estrategia de Campana completa: analisis, táctica y ejecucion.",
       MH, Inches(4.8), SAFW * 0.68, Emu(550000),
       fnt=FONT_B, sz=Pt(14), color=TEXT_SOFT)

    # Date chip bottom-left
    tx(s, "NOVIEMBRE 2026",
       MH, SH - MV - Emu(380000), Inches(3), Emu(320000),
       fnt=FONT_M, sz=Pt(10), color=TEXT_DIM)

    # Bottom accent line
    bar(s, Inches(0), SH - S48, SW, S4)
    return s


def s02_resumen(prs):
    s = blank(prs)
    bg(s, BG_DARK)

    pill(s, "01  RESUMEN EJECUTIVO", MH, MV)

    tx(s, "En 7 dias podemos cerrar el equivalente a 2 meses de trabajo.",
       MH, MV + Emu(460000), SAFW, Emu(950000),
       fnt=FONT_H, sz=Pt(28), bold=True, color=TEXT_ON_DARK)

    # 4 metric cards, 2x2 grid
    metrics = [
        ("28 NOV",       "Fecha del evento"),
        ("3 CUPOS",      "Capacidad maxima"),
        ("$1.200.000",   "Objetivo de facturacion"),
        ("7 DIAS",       "Ventana de campana"),
    ]

    card_w = (SAFW - S24) / 2
    card_h = Emu(1200000)

    for i, (val, lbl) in enumerate(metrics):
        col = i % 2
        row = i // 2
        cx = MH + col * (card_w + S24)
        cy = MV + Emu(1280000) + row * (card_h + S16)

        rct(s, cx, cy, card_w, card_h, SURFACE_2)
        bar(s, cx, cy, S8, card_h)

        tx(s, val,
           cx + S16, cy + S12, card_w - S32, Emu(650000),
           fnt=FONT_M, sz=Pt(24), bold=True, color=ACCENT)

        tx(s, lbl,
           cx + S16, cy + Emu(680000), card_w - S32, Emu(380000),
           fnt=FONT_B, sz=Pt(11), color=TEXT_SOFT)

    return s


def s03_div_mercado(prs):
    return divider_slide(prs, "01", "ANALISIS", "DE MERCADO")


def s04_oportunidad(prs):
    s = blank(prs)
    bg(s, BG_LIGHT)

    tx(s, "El Q4 es el pico de decision digital del dueno de PyME.",
       MH, MV, SAFW, Inches(1.0),
       fnt=FONT_H, sz=Pt(26), bold=True, color=TEXT_PRIMARY)

    bar(s, MH, MV + Inches(1.1), Inches(2.2), S4)

    datos = [
        "+180%  en busquedas de 'sitio web para mi negocio' cada noviembre (Google Argentina).",
        "7 de cada 10  duenos de PyME considera invertir en digital en Q4 segun Mercado Ads.",
        "3 de cada 10  PyMEs argentinas tiene web propia — la oportunidad sigue siendo enorme.",
        "Agencias en Q4  suben precios y alargan plazos. Nosotros arrancamos en 2-3 semanas.",
    ]

    for i, dato in enumerate(datos):
        y = MV + Inches(1.4) + i * Emu(700000)
        bar(s, MH, y + Emu(180000), S8, S8)
        tx(s, dato,
           MH + S24, y, SAFW - S24, Emu(620000),
           fnt=FONT_B, sz=Pt(13), color=TEXT_PRIMARY)

    return s


def s05_competencia(prs):
    s = blank(prs)
    bg(s, BG_LIGHT)

    # Dark left panel
    lw = SW * 0.44
    rct(s, Inches(0), Inches(0), lw, SH, BG_DARK)

    # Left: agencias
    tx(s, "AGENCIAS",
       S32, MV, lw - S64, Emu(360000),
       fnt=FONT_M, sz=Pt(10), color=TEXT_DIM)

    tx(s, "El costo\noculto.",
       S32, MV + Emu(480000), lw - S64, Emu(1100000),
       fnt=FONT_H, sz=Pt(36), bold=True, color=TEXT_ON_DARK)

    contras = [
        "Desde $1.500.000 ARS promedio",
        "60 a 90 dias de desarrollo",
        "Cuotas mensuales sin fecha de fin",
        "Cambian el equipo a mitad de proyecto",
        "PM intermediario: nada directo",
    ]
    for i, c in enumerate(contras):
        y = Inches(2.6) + i * Emu(560000)
        tx(s, "x  " + c,
           S32, y, lw - S64, Emu(490000),
           fnt=FONT_B, sz=Pt(12), color=TEXT_SOFT)

    # Right: Selva Digital
    rx = lw + S48
    rw = SW - lw - S48 - MH

    tx(s, "SELVA DIGITAL",
       rx, MV, rw, Emu(360000),
       fnt=FONT_M, sz=Pt(10), color=ACCENT)

    tx(s, "Pago unico.\nSin vueltas.",
       rx, MV + Emu(480000), rw, Emu(1100000),
       fnt=FONT_H, sz=Pt(36), bold=True, color=TEXT_PRIMARY)

    pros = [
        "$250.000 – $700.000 segun el proyecto",
        "2 a 3 semanas para un sitio completo",
        "Pago unico, sin cuotas mensuales",
        "Hablás con Christian de principio a fin",
        "Arranca solo cuando vos aprobas el diseno",
    ]
    for i, p in enumerate(pros):
        y = Inches(2.6) + i * Emu(560000)
        tx(s, "v  " + p,
           rx, y, rw, Emu(490000),
           fnt=FONT_B, sz=Pt(12), bold=(i < 2), color=TEXT_PRIMARY)

    return s


def s06_div_estrategia(prs):
    return divider_slide(prs, "02", "ESTRATEGIA")


def s07_estrategia(prs):
    s = blank(prs)
    bg(s, BG_LIGHT)

    tx(s, "Una oferta. Una fecha. Tres cupos.\nLa escasez real genera accion.",
       MH, MV, SAFW, Inches(1.5),
       fnt=FONT_H, sz=Pt(28), bold=True, color=TEXT_PRIMARY)

    bar(s, MH, MV + Inches(1.6), Inches(2.8), S4)

    tx(s, "No es un descuento de precio: es una garantia de atencion exclusiva. "
          "El dueno de PyME no necesita un precio mas bajo — necesita certeza de "
          "que su proyecto arranca esta semana y termina en tiempo y forma.",
       MH, MV + Inches(1.9), SAFW * 0.70, Inches(1.6),
       fnt=FONT_B, sz=Pt(14), color=TEXT_SOFT)

    # Servicios disponibles
    tx(s, "LA OFERTA BLACK FRIDAY:",
       MH, MV + Inches(3.7), SAFW, Emu(360000),
       fnt=FONT_M, sz=Pt(10), color=TEXT_DIM)

    ofertas = [
        "LANDING PAGE  $250.000",
        "SITIO WEB  $400.000",
        "E-COMMERCE  $700.000",
    ]
    px = MH
    for oferta in ofertas:
        p_shape = pill(s, oferta, px, MV + Inches(4.2))
        # advance x by pill width
        px += Emu(int(len(oferta) * 85000 + 400000))

    tx(s, "Sena del 40% confirma el cupo. Solo hasta completar los 3 proyectos.",
       MH, MV + Inches(5.1), SAFW, Emu(360000),
       fnt=FONT_B, sz=Pt(11), italic=True, color=TEXT_DIM)

    return s


def s08_posicionamiento(prs):
    s = blank(prs)
    bg(s, SURFACE_1)

    # Big decorative quotes
    tx(s, "“",
       MH - S8, Inches(0.3), Emu(1000000), Emu(1600000),
       fnt=FONT_H, sz=Pt(140), bold=True, color=ACCENT)

    tx(s, "Mientras la competencia\ndescuenta el precio,\nnosotros descontamos\nel riesgo.",
       MH, Inches(1.6), SAFW * 0.80, Inches(3.4),
       fnt=FONT_H, sz=Pt(30), italic=True, color=TEXT_ON_DARK)

    tx(s, "Posicionamiento Selva Digital  Black Friday 2026",
       MH, Inches(5.3), SAFW, Emu(360000),
       fnt=FONT_B, sz=Pt(11), color=TEXT_DIM)

    return s


def s09_div_tacticas(prs):
    return divider_slide(prs, "03", "TACTICAS")


def s10_precampana(prs):
    s = blank(prs)
    bg(s, BG_LIGHT)

    pill(s, "SEMANA PREVIA  24 - 27 NOV", MH, MV)

    tx(s, "Calentar la audiencia el lunes duplica la conversion del viernes.",
       MH, MV + Emu(450000), SAFW, Emu(850000),
       fnt=FONT_H, sz=Pt(24), bold=True, color=TEXT_PRIMARY)

    acciones = [
        ("LUN 24/11", "Post carrusel en Instagram: '¿Tu negocio tiene web?' — problema > solucion > resultado real del portfolio."),
        ("MAR 25/11", "Stories con countdown 72 hs + early access por WhatsApp a lista de leads anteriores."),
        ("MIE 26/11", "Mensaje 1:1 a contactos warm: 'Te reservo un cupo si me avísas antes del viernes'."),
        ("JUE 27/11", "Reels antes/despues con caso real: El Fogon (x2.3 ticket) o MegaMuebles (+34% leads)."),
    ]

    for i, (fecha, accion) in enumerate(acciones):
        y = MV + Inches(1.7) + i * Emu(710000)
        rct(s, MH, y, Inches(1.1), Emu(590000), ACCENT_TINT, lc=ACCENT)
        tx(s, fecha,
           MH + S8, y + S8, Inches(1.0), Emu(500000),
           fnt=FONT_M, sz=Pt(9), bold=True, color=ACCENT)
        tx(s, accion,
           MH + Inches(1.25), y + S12, SAFW - Inches(1.25), Emu(560000),
           fnt=FONT_B, sz=Pt(12), color=TEXT_PRIMARY)

    tx(s, "Horario pico de publicacion: 19 - 21 hs (Argentina)",
       MH, SH - MV - Emu(380000), SAFW, Emu(320000),
       fnt=FONT_M, sz=Pt(10), color=TEXT_DIM)

    return s


def s11_bf_day(prs):
    s = blank(prs)
    bg(s, BG_DARK)

    pill(s, "DIA 0  VIERNES 28 NOV", MH, MV)

    tx(s, "Urgencia visible + prueba social en tiempo real.",
       MH, MV + Emu(450000), SAFW, Emu(800000),
       fnt=FONT_H, sz=Pt(26), bold=True, color=TEXT_ON_DARK)

    horas = [
        ("09:00",  "Publicar oferta completa en IG feed + primera story con link en bio."),
        ("10:00",  "Blast por WhatsApp a lista warm (maximo 50 contactos seleccionados)."),
        ("14:00",  "Story update: 'CUPO 1 TOMADO — quedan 2.'"),
        ("20:00",  "Story update: 'CUPO 2 TOMADO — queda 1 · manana cerramos.'"),
    ]

    card_w = (SAFW - S24) / 2
    card_h = Emu(850000)

    for i, (hora, accion) in enumerate(horas):
        col = i % 2
        row = i // 2
        x = MH + col * (card_w + S24)
        y = MV + Inches(1.55) + row * (card_h + S12)

        rct(s, x, y, card_w, card_h, SURFACE_2)
        bar(s, x, y, S8, card_h)

        tx(s, hora,
           x + S16, y + S12, card_w - S32, Emu(400000),
           fnt=FONT_M, sz=Pt(18), bold=True, color=ACCENT)

        tx(s, accion,
           x + S16, y + Emu(450000), card_w - S32, Emu(380000),
           fnt=FONT_B, sz=Pt(11), color=TEXT_SOFT)

    tx(s, "Regla de oro: responder toda consulta en menos de 30 minutos durante el dia.",
       MH, SH - MV - Emu(360000), SAFW, Emu(320000),
       fnt=FONT_B, sz=Pt(11), italic=True, color=TEXT_DIM)

    return s


def s12_cyber(prs):
    s = blank(prs)
    bg(s, BG_LIGHT)

    pill(s, "SEGUIMIENTO  LUN - MAR 1 - 2 DIC", MH, MV)

    tx(s, "Cyber Monday convierte a los que dudaron el viernes.",
       MH, MV + Emu(450000), SAFW, Emu(800000),
       fnt=FONT_H, sz=Pt(26), bold=True, color=TEXT_PRIMARY)

    bar(s, MH, MV + Inches(1.5), Inches(1.8), S4)

    acciones = [
        "Lunes 1/12 a las 9 AM: 'Ultima oportunidad — queda 1 cupo' (solo si es cierto).",
        "Mensaje distinto al del viernes: 'Ya arrancamos con 2 proyectos. Te sumo si me confirmas hoy.'",
        "Mostrar avance en stories si ya hay un cliente firmado (genera prueba social en vivo).",
        "Cerrar campana el martes 2/12 al mediodia con story de cierre: 'Cupos agotados. Proxima ventana: enero 2027.'",
        "Guardar lista de interesados que no cerraron — son leads para la proxima ventana.",
    ]

    for i, accion in enumerate(acciones):
        y = MV + Inches(1.8) + i * Emu(650000)
        bar(s, MH, y + Emu(180000), S8, S8)
        tx(s, accion,
           MH + S24, y, SAFW - S24, Emu(580000),
           fnt=FONT_B, sz=Pt(12), color=TEXT_PRIMARY)

    return s


def s13_canales(prs):
    s = blank(prs)
    bg(s, BG_LIGHT)

    tx(s, "Tres canales, un solo mensaje: tu negocio necesita web.",
       MH, MV, SAFW, Emu(800000),
       fnt=FONT_H, sz=Pt(24), bold=True, color=TEXT_PRIMARY)

    canales = [
        ("INSTAGRAM",  "Alcance organico",   "Carrusel + Stories + Reels\nCold → warm audience"),
        ("WHATSAPP",   "Cierre 1 a 1",        "Lista de leads warm\nConversion directa"),
        ("DM + EMAIL", "Nurturing",           "Leads anteriores\nProspectos indecisos"),
    ]

    n = len(canales)
    cw = (SAFW - (n - 1) * S24) / n
    ch = Inches(4.0)
    cy = MV + Emu(950000)

    for i, (nombre, tipo, desc) in enumerate(canales):
        cx = MH + i * (cw + S24)
        rct(s, cx, cy, cw, ch, ACCENT_TINT, lc=ACCENT, lw=Emu(9525))
        bar(s, cx, cy, cw, S4)

        tx(s, nombre,
           cx + S16, cy + S16, cw - S32, Emu(420000),
           fnt=FONT_M, sz=Pt(11), bold=True, color=ACCENT)

        tx(s, tipo,
           cx + S16, cy + Emu(500000), cw - S32, Emu(400000),
           fnt=FONT_H, sz=Pt(15), bold=True, color=TEXT_PRIMARY)

        tx(s, desc,
           cx + S16, cy + Emu(980000), cw - S32, Emu(950000),
           fnt=FONT_B, sz=Pt(12), color=TEXT_SOFT)

    return s


def s14_div_pasos(prs):
    return divider_slide(prs, "04", "PROXIMOS", "PASOS")


def s15_timeline(prs):
    s = blank(prs)
    bg(s, BG_LIGHT)

    tx(s, "De hoy a noviembre: 4 entregas, 0 improvisaciones.",
       MH, MV, SAFW, Emu(800000),
       fnt=FONT_H, sz=Pt(24), bold=True, color=TEXT_PRIMARY)

    pasos = [
        ("26 MAY",    "Brief + identidad",     "Definir oferta, mensajes y piezas necesarias"),
        ("JUNIO",     "Copy + diseno",          "Carrusel, stories, reels y copy de WA"),
        ("OCTUBRE",   "Warmup de audiencia",    "Early access, countdown, nurturing leads"),
        ("28 NOV",    "Black Friday LIVE",      "Ejecucion en tiempo real, cierre de cupos"),
    ]

    n = len(pasos)
    step_w = SAFW / n
    badge_h = Emu(540000)
    badge_y = Inches(2.0)

    for i, (fecha, titulo, detalle) in enumerate(pasos):
        x = MH + step_w * i
        bx = x + (step_w - Emu(720000)) / 2

        # Connector to next
        if i < n - 1:
            bar(s, x + step_w * 0.58, badge_y + badge_h / 2 - S4 / 2,
                step_w * 0.38, Emu(18000))

        # Date badge
        badge_fill = ACCENT if i == n - 1 else SURFACE_1
        date_color = CTA_TEXT if i == n - 1 else ACCENT
        rct(s, bx, badge_y, Emu(720000), badge_h, badge_fill)

        tx(s, fecha,
           bx, badge_y + S12, Emu(720000), badge_h - S24,
           fnt=FONT_M, sz=Pt(10), bold=True, color=date_color,
           align=PP_ALIGN.CENTER)

        tx(s, titulo,
           x + S8, badge_y + badge_h + S16, step_w - S16, Emu(500000),
           fnt=FONT_H, sz=Pt(13), bold=True, color=TEXT_PRIMARY,
           align=PP_ALIGN.CENTER)

        tx(s, detalle,
           x + S8, badge_y + badge_h + Emu(600000), step_w - S16, Emu(650000),
           fnt=FONT_B, sz=Pt(11), color=TEXT_SOFT,
           align=PP_ALIGN.CENTER)

    tx(s, "Estado: lista para arrancar — solo falta la aprobacion.",
       MH, SH - MV - Emu(360000), SAFW, Emu(320000),
       fnt=FONT_B, sz=Pt(11), italic=True, color=TEXT_DIM)

    return s


def s16_cta(prs):
    s = blank(prs)
    bg(s, BG_DARK)

    tx(s, "El mejor momento para\npreparar Black Friday\nfue ayer.",
       MH, Inches(1.2), SAFW * 0.78, Inches(2.8),
       fnt=FONT_H, sz=Pt(44), bold=True, color=TEXT_ON_DARK)

    tx(s, "El segundo mejor es hoy.",
       MH, Inches(3.9), SAFW * 0.78, Emu(900000),
       fnt=FONT_H, sz=Pt(32), bold=True, color=ACCENT)

    pill(s, "+54 9 3548 550334",        MH, Inches(5.2))
    pill(s, "info.selvadigital@gmail.com", MH + Inches(3.1), Inches(5.2))
    pill(s, "INICIO CAMPANA: 26 MAYO 2026", MH, Inches(5.9))

    bar(s, Inches(0), SH - S48, SW, S4)
    return s


# ============================================================
# BUILD
# ============================================================

def build(output_path):
    prs = new_prs()

    builders = [
        s01_cover,
        s02_resumen,
        s03_div_mercado,
        s04_oportunidad,
        s05_competencia,
        s06_div_estrategia,
        s07_estrategia,
        s08_posicionamiento,
        s09_div_tacticas,
        s10_precampana,
        s11_bf_day,
        s12_cyber,
        s13_canales,
        s14_div_pasos,
        s15_timeline,
        s16_cta,
    ]

    for fn in builders:
        fn(prs)

    out = os.path.join(os.path.dirname(__file__), output_path)
    os.makedirs(os.path.dirname(os.path.abspath(out)), exist_ok=True)
    prs.save(out)
    print(f"OK  {len(prs.slides)} slides guardados en: {out}")


if __name__ == "__main__":
    build("campana-black-friday-2026-06.pptx")
